#!/usr/bin/env python3
"""
문진 폴더의 docx, pptx 파일 내용을 추출하는 스크립트
"""

import os
import sys
from pathlib import Path
from docx import Document
from pptx import Presentation

# 문진 폴더 경로
FOLDER_PATH = Path("/Users/vibetj/coding/center/docs/문진")
OUTPUT_DIR = Path("/Users/vibetj/coding/center/docs/문진/extracted_texts")

def extract_docx(file_path):
    """docx 파일에서 텍스트 추출"""
    try:
        doc = Document(file_path)
        content = []
        
        # 문단 추출
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                content.append(paragraph.text)
        
        # 표 추출
        for table in doc.tables:
            content.append("\n[표]")
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    content.append(" | ".join(row_text))
            content.append("[표 끝]\n")
        
        return "\n".join(content)
    except Exception as e:
        return f"오류 발생: {str(e)}"

def extract_pptx(file_path):
    """pptx 파일에서 텍스트 추출"""
    try:
        prs = Presentation(file_path)
        content = []
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            content.append(f"\n=== 슬라이드 {slide_idx} ===")
            
            # 슬라이드의 모든 shape 확인
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text)
                
                # 표가 있는 경우
                if shape.has_table:
                    content.append("\n[표]")
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            content.append(" | ".join(row_text))
                    content.append("[표 끝]\n")
        
        return "\n".join(content)
    except Exception as e:
        return f"오류 발생: {str(e)}"

def main():
    # 출력 디렉토리 생성
    OUTPUT_DIR.mkdir(exist_ok=True)
    
    # docx 파일 처리
    docx_files = list(FOLDER_PATH.glob("*.docx"))
    
    print(f"\n발견된 docx 파일: {len(docx_files)}개")
    
    for file_path in docx_files:
        print(f"\n처리 중: {file_path.name}")
        content = extract_docx(file_path)
        
        # 출력 파일명 생성
        output_file = OUTPUT_DIR / f"{file_path.stem}.txt"
        
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(f"=== {file_path.name} 내용 추출 ===\n\n")
            f.write(content)
        
        print(f"저장됨: {output_file}")
    
    # pptx 파일 처리
    pptx_files = list(FOLDER_PATH.glob("*.pptx"))
    
    print(f"\n발견된 pptx 파일: {len(pptx_files)}개")
    
    for file_path in pptx_files:
        print(f"\n처리 중: {file_path.name}")
        content = extract_pptx(file_path)
        
        # 출력 파일명 생성
        output_file = OUTPUT_DIR / f"{file_path.stem}.txt"
        
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(f"=== {file_path.name} 내용 추출 ===\n\n")
            f.write(content)
        
        print(f"저장됨: {output_file}")

if __name__ == "__main__":
    main()